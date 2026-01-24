"""
Multimodal document processing for RAG.

Supports:
- Text extraction from various formats (PDF, DOCX, PPTX, CSV, TXT)
- OCR for scanned documents and images
- Structured element extraction (titles, paragraphs, tables, images)
- Title-based chunking
- Image-to-text conversion via vision models
"""

import os
import io
import logging
from typing import List, Tuple, Optional, Dict, Any
from enum import Enum
from pathlib import Path
import hashlib

logger = logging.getLogger(__name__)

# File type detection
class DocumentType(Enum):
    """Supported document types"""
    PDF = "pdf"
    DOCX = "docx"
    PPTX = "pptx"
    CSV = "csv"
    TXT = "txt"
    IMAGE = "image"
    UNKNOWN = "unknown"


class ProcessedElement:
    """Represents a single extracted element from a document"""
    
    def __init__(
        self,
        element_type: str,
        content: str,
        metadata: Optional[Dict[str, Any]] = None,
        images: Optional[List[bytes]] = None,
    ):
        """
        Args:
            element_type: Type of element (text, title, table, image, etc.)
            content: Text or HTML content
            metadata: Optional metadata dict
            images: Optional list of image bytes
        """
        self.element_type = element_type
        self.content = content
        self.metadata = metadata or {}
        self.images = images or []
    
    def to_dict(self) -> Dict[str, Any]:
        """Convert to dict representation"""
        return {
            "type": self.element_type,
            "content": self.content,
            "metadata": self.metadata,
            "has_images": len(self.images) > 0,
        }


def detect_file_type(filename: str, file_bytes: Optional[bytes] = None) -> DocumentType:
    """
    Detect document type by extension and optionally by file signature.
    
    Args:
        filename: Filename with extension
        file_bytes: Optional file bytes for magic number detection
    
    Returns:
        DocumentType enum
    """
    ext = Path(filename).suffix.lower()
    
    # Extension-based detection
    ext_map = {
        ".pdf": DocumentType.PDF,
        ".docx": DocumentType.DOCX,
        ".doc": DocumentType.DOCX,
        ".pptx": DocumentType.PPTX,
        ".ppt": DocumentType.PPTX,
        ".csv": DocumentType.CSV,
        ".txt": DocumentType.TXT,
        ".png": DocumentType.IMAGE,
        ".jpg": DocumentType.IMAGE,
        ".jpeg": DocumentType.IMAGE,
        ".gif": DocumentType.IMAGE,
        ".webp": DocumentType.IMAGE,
    }
    
    if ext in ext_map:
        return ext_map[ext]
    
    # Magic number detection if bytes provided
    if file_bytes:
        if file_bytes.startswith(b"%PDF"):
            return DocumentType.PDF
        if file_bytes.startswith(b"PK"):  # ZIP-based (DOCX, PPTX, XLSX)
            # Read as ZIP to determine type
            try:
                import zipfile
                with zipfile.ZipFile(io.BytesIO(file_bytes)) as zf:
                    names = zf.namelist()
                    if any("word/" in n for n in names):
                        return DocumentType.DOCX
                    if any("ppt/" in n for n in names):
                        return DocumentType.PPTX
            except Exception:
                pass
        if file_bytes.startswith(b"\xff\xd8\xff"):  # JPEG
            return DocumentType.IMAGE
        if file_bytes.startswith(b"\x89PNG"):  # PNG
            return DocumentType.IMAGE
    
    return DocumentType.UNKNOWN


async def extract_text_from_pdf(file_bytes: bytes, ocr_enabled: bool = True) -> List[ProcessedElement]:
    """
    Extract text from PDF, with automatic OCR for scanned documents.
    
    Args:
        file_bytes: PDF file bytes
        ocr_enabled: Enable OCR for scanned PDFs
    
    Returns:
        List of ProcessedElement objects
    """
    elements = []
    
    try:
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(file_bytes))
        
        for page_num, page in enumerate(reader.pages):
            try:
                text = page.extract_text()
                
                # Check if page is mostly scanned (little text extracted)
                if ocr_enabled and (not text or len(text.strip()) < 100):
                    logger.info(f"Page {page_num + 1} appears scanned, attempting OCR")
                    # Fallback to OCR
                    try:
                        text = await _ocr_page(file_bytes, page_num)
                    except Exception as e:
                        logger.warning(f"OCR failed for page {page_num + 1}: {e}")
                        text = text or "[OCR failed - page content unavailable]"
                
                if text and text.strip():
                    elements.append(ProcessedElement(
                        element_type="text",
                        content=text.strip(),
                        metadata={"page": page_num + 1, "source": "pdf_text_extraction"},
                    ))
            except Exception as e:
                logger.warning(f"Error extracting text from PDF page {page_num + 1}: {e}")
        
        return elements
    
    except Exception as e:
        logger.error(f"Error reading PDF: {e}")
        raise


async def extract_text_from_docx(file_bytes: bytes) -> List[ProcessedElement]:
    """Extract text from DOCX file"""
    elements = []
    
    try:
        from docx import Document
        doc = Document(io.BytesIO(file_bytes))
        
        for para in doc.paragraphs:
            if para.text.strip():
                elements.append(ProcessedElement(
                    element_type="paragraph",
                    content=para.text.strip(),
                    metadata={"source": "docx"},
                ))
        
        # Extract tables
        for table in doc.tables:
            table_text = "\n".join([
                "\t".join(cell.text for cell in row.cells)
                for row in table.rows
            ])
            if table_text.strip():
                elements.append(ProcessedElement(
                    element_type="table",
                    content=table_text,
                    metadata={"source": "docx_table"},
                ))
        
        return elements
    
    except Exception as e:
        logger.error(f"Error reading DOCX: {e}")
        raise


async def extract_text_from_pptx(file_bytes: bytes) -> List[ProcessedElement]:
    """Extract text from PPTX file"""
    elements = []
    
    try:
        from pptx import Presentation
        prs = Presentation(io.BytesIO(file_bytes))
        
        for slide_num, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    elements.append(ProcessedElement(
                        element_type="text",
                        content=shape.text.strip(),
                        metadata={
                            "slide": slide_num + 1,
                            "source": "pptx",
                        },
                    ))
                
                # Extract tables from shapes
                if shape.has_table:
                    table = shape.table
                    table_text = "\n".join([
                        "\t".join(cell.text for cell in row.cells)
                        for row in table.rows
                    ])
                    if table_text.strip():
                        elements.append(ProcessedElement(
                            element_type="table",
                            content=table_text,
                            metadata={
                                "slide": slide_num + 1,
                                "source": "pptx_table",
                            },
                        ))
        
        return elements
    
    except Exception as e:
        logger.error(f"Error reading PPTX: {e}")
        raise


async def extract_text_from_csv(file_bytes: bytes) -> List[ProcessedElement]:
    """Extract text from CSV file"""
    elements = []
    
    try:
        import csv
        content = file_bytes.decode("utf-8", errors="replace")
        reader = csv.reader(io.StringIO(content))
        
        # Read all rows
        rows = list(reader)
        
        if not rows:
            return elements
        
        # Add header
        if rows:
            header_text = " | ".join(str(cell) for cell in rows[0])
            elements.append(ProcessedElement(
                element_type="table_header",
                content=header_text,
                metadata={"source": "csv"},
            ))
        
        # Add data rows as text
        for row_num, row in enumerate(rows[1:]):
            row_text = " | ".join(str(cell) for cell in row)
            if row_text.strip():
                elements.append(ProcessedElement(
                    element_type="table_row",
                    content=row_text,
                    metadata={
                        "row": row_num + 1,
                        "source": "csv",
                    },
                ))
        
        return elements
    
    except Exception as e:
        logger.error(f"Error reading CSV: {e}")
        raise


async def extract_text_from_image(file_bytes: bytes, use_ocr: bool = True) -> List[ProcessedElement]:
    """
    Extract text from image file.
    
    Args:
        file_bytes: Image file bytes
        use_ocr: Enable OCR extraction
    
    Returns:
        List of ProcessedElement objects
    """
    elements = []
    
    if not use_ocr:
        return [ProcessedElement(
            element_type="image",
            content="[Image - OCR not enabled]",
            metadata={"source": "image"},
            images=[file_bytes],
        )]
    
    try:
        text = await _ocr_image(file_bytes)
        elements.append(ProcessedElement(
            element_type="image_text",
            content=text,
            metadata={"source": "image_ocr"},
            images=[file_bytes],
        ))
    except Exception as e:
        logger.warning(f"OCR failed for image: {e}")
        elements.append(ProcessedElement(
            element_type="image",
            content="[Image - OCR processing failed]",
            metadata={"source": "image", "ocr_error": str(e)},
            images=[file_bytes],
        ))
    
    return elements


async def extract_text_from_txt(file_bytes: bytes) -> List[ProcessedElement]:
    """Extract text from plain text file"""
    try:
        content = file_bytes.decode("utf-8", errors="replace").strip()
        
        if not content:
            return []
        
        return [ProcessedElement(
            element_type="text",
            content=content,
            metadata={"source": "txt"},
        )]
    except Exception as e:
        logger.error(f"Error reading TXT: {e}")
        raise


async def extract_elements_from_file(
    filename: str,
    file_bytes: bytes,
    doc_type: Optional[DocumentType] = None,
) -> List[ProcessedElement]:
    """
    Extract structured elements from any supported file type.
    
    Args:
        filename: Original filename
        file_bytes: File contents
        doc_type: Optional pre-detected document type
    
    Returns:
        List of ProcessedElement objects
    """
    # Auto-detect type if not provided
    if doc_type is None:
        doc_type = detect_file_type(filename, file_bytes)
    
    logger.info(f"Extracting elements from {filename} (detected type: {doc_type.value})")
    
    extractors = {
        DocumentType.PDF: extract_text_from_pdf,
        DocumentType.DOCX: extract_text_from_docx,
        DocumentType.PPTX: extract_text_from_pptx,
        DocumentType.CSV: extract_text_from_csv,
        DocumentType.IMAGE: extract_text_from_image,
        DocumentType.TXT: extract_text_from_txt,
    }
    
    extractor = extractors.get(doc_type)
    if not extractor:
        raise ValueError(f"Unsupported document type: {doc_type.value}")
    
    elements = await extractor(file_bytes)
    logger.info(f"Extracted {len(elements)} elements from {filename}")
    
    return elements


def chunk_elements_by_title(
    elements: List[ProcessedElement],
    max_chunk_chars: int = 1200,
    merge_threshold_chars: int = 300,
) -> List[str]:
    """
    Chunk elements intelligently by title/section boundaries.
    
    This combines multiple elements into meaningful chunks based on:
    - Title elements act as chunk boundaries
    - Similar-level content is grouped together
    - Chunks are limited by character count
    
    Args:
        elements: List of ProcessedElement objects
        max_chunk_chars: Hard limit for chunk size
        merge_threshold_chars: Merge small chunks smaller than this
    
    Returns:
        List of text chunks ready for embedding
    """
    chunks = []
    current_chunk_parts = []
    current_chunk_length = 0
    
    for element in elements:
        content = element.content
        element_length = len(content)
        
        # Titles and section headers should start new chunks
        is_boundary = element.element_type in ["title", "heading", "section"]
        
        # If current chunk is not empty and adding this would exceed max, flush it
        if current_chunk_parts and current_chunk_length + element_length > max_chunk_chars:
            chunk_text = "\n\n".join(current_chunk_parts)
            if len(chunk_text) >= merge_threshold_chars or not chunks:
                chunks.append(chunk_text)
            current_chunk_parts = []
            current_chunk_length = 0
        
        # If this element is too large, split it
        if element_length > max_chunk_chars:
            # Flush current chunk first
            if current_chunk_parts:
                chunk_text = "\n\n".join(current_chunk_parts)
                if len(chunk_text) >= merge_threshold_chars or not chunks:
                    chunks.append(chunk_text)
                current_chunk_parts = []
                current_chunk_length = 0
            
            # Split large element into smaller pieces
            for i in range(0, element_length, max_chunk_chars):
                chunk_part = content[i:i + max_chunk_chars]
                if chunk_part.strip():
                    chunks.append(chunk_part.strip())
        else:
            # Add to current chunk
            current_chunk_parts.append(content)
            current_chunk_length += element_length + 2  # +2 for separator
            
            # Start new chunk at boundaries if chunk is reasonable size
            if is_boundary and current_chunk_length > 500:
                chunk_text = "\n\n".join(current_chunk_parts)
                if len(chunk_text) >= merge_threshold_chars:
                    chunks.append(chunk_text)
                current_chunk_parts = []
                current_chunk_length = 0
    
    # Add final chunk
    if current_chunk_parts:
        chunk_text = "\n\n".join(current_chunk_parts)
        if len(chunk_text) >= merge_threshold_chars or not chunks:
            chunks.append(chunk_text)
    
    logger.info(f"Created {len(chunks)} chunks from {len(elements)} elements")
    return [c.strip() for c in chunks if c.strip()]


async def image_to_text_summary(
    image_bytes: bytes,
    use_vision_api: bool = True,
    fallback_text: str = "",
) -> str:
    """
    Convert image to text using OpenAI Vision API (GPT-4 Turbo).
    
    Args:
        image_bytes: Image file bytes
        use_vision_api: Whether to use vision API
        fallback_text: Fallback text if vision API is not available
    
    Returns:
        Text description/summary of image
    """
    if not use_vision_api or not os.getenv("OPENAI_API_KEY"):
        return fallback_text or "[Image - Vision API not configured]"
    
    try:
        from openai import OpenAI
        import base64
        
        client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
        
        # Encode image as base64
        image_b64 = base64.standard_b64encode(image_bytes).decode("utf-8")
        
        # Call vision API with GPT-4 Turbo
        response = client.chat.completions.create(
            model="gpt-4-turbo",
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": (
                                "Describe this image in detail for document indexing. "
                                "Include: main subjects, text visible in image, charts/diagrams, "
                                "colors, layout, and any important details. Be comprehensive."
                            ),
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:image/jpeg;base64,{image_b64}"
                            },
                        },
                    ],
                }
            ],
            max_tokens=500,
        )
        
        return response.choices[0].message.content or "[Image - No description generated]"
    
    except Exception as e:
        logger.warning(f"Vision API failed: {e}")
        return fallback_text or f"[Image - Vision API processing failed: {str(e)}]"


async def _ocr_image(image_bytes: bytes) -> str:
    """
    Run OCR on image using PaddleOCR (fallback to Tesseract).
    Returns placeholder if both fail - OCR is optional.
    
    Args:
        image_bytes: Image file bytes
    
    Returns:
        Extracted text or placeholder
    """
    # Try PaddleOCR first (more accurate, supports more languages)
    try:
        from paddleocr import PaddleOCR
        ocr = PaddleOCR(use_angle_cls=True, lang="en")
        
        result = ocr.ocr(image_bytes, cls=True)
        texts = []
        for line in result:
            if line:
                for word_info in line:
                    texts.append(word_info[1][0])
        
        if texts:
            logger.info(f"PaddleOCR extracted {len(texts)} text elements")
            return " ".join(texts)
        else:
            logger.debug("PaddleOCR: No text detected, trying Tesseract")
    
    except Exception as e:
        logger.debug(f"PaddleOCR failed: {e}, trying Tesseract")
    
    # Fallback to Tesseract
    try:
        import pytesseract
        from PIL import Image
        
        img = Image.open(io.BytesIO(image_bytes))
        text = pytesseract.image_to_string(img)
        
        if text.strip():
            logger.info(f"Tesseract extracted text")
            return text.strip()
        else:
            logger.debug("Tesseract: No text detected")
    
    except Exception as e:
        logger.debug(f"Tesseract OCR failed: {e}")
    
    # If both fail, return placeholder - OCR is optional
    logger.debug("Both OCR engines failed - using placeholder")
    return "[Image content - OCR unavailable]"


async def _ocr_tesseract(image_bytes: bytes) -> str:
    """Fallback OCR using Tesseract - returns placeholder on failure"""
    try:
        import pytesseract
        from PIL import Image
        
        img = Image.open(io.BytesIO(image_bytes))
        text = pytesseract.image_to_string(img)
        
        return text.strip() if text.strip() else "[OCR: No text detected]"
    
    except Exception as e:
        logger.debug(f"Tesseract OCR failed: {e}")
        return "[Tesseract unavailable]"


async def _ocr_page(file_bytes: bytes, page_num: int) -> str:
    """Convert PDF page to image and run OCR - gracefully handles failures"""
    try:
        # Convert page to image using PyMuPDF
        import fitz  # PyMuPDF
        
        doc = fitz.open(stream=io.BytesIO(file_bytes), filetype="pdf")
        pix = doc[page_num].get_pixmap(matrix=fitz.Matrix(2, 2))
        image_bytes = pix.tobytes("png")
        
        return await _ocr_image(image_bytes)
    
    except Exception as e:
        logger.debug(f"PDF to image OCR failed: {e}")
        return "[Page image - OCR unavailable]"


def compute_file_hash(file_bytes: bytes) -> str:
    """Compute SHA-256 hash of file for deduplication"""
    return hashlib.sha256(file_bytes).hexdigest()
